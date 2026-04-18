"""Local file content extraction — parse local files into structured sections.

Supports .md, .txt (built-in), and via MarkItDown: .docx, .pptx, .doc, .html,
.mhtml, .pdf, .xlsx, .xls, .csv, .json, .xml, .epub, .msg, .zip.
Legacy extractors (python-docx, python-pptx, mammoth, trafilatura) serve as
fallbacks when MarkItDown fails or is disabled.
"""

import email
import email.policy
import pathlib
import re
from datetime import datetime
from typing import List, Optional, Tuple, TYPE_CHECKING

from .exceptions import LocalFileError
from .models import ArticleInfo, ArticleSection, ExtractedImage

if TYPE_CHECKING:
    from .config import LocalFilesConfig


# ---------------------------------------------------------------------------
# Dispatcher
# ---------------------------------------------------------------------------

def extract_local_file(
    file_path: str,
    config: "LocalFilesConfig",
    extract_images: bool = False,
    markitdown_config=None,
) -> Tuple[ArticleInfo, List[ArticleSection], List[ExtractedImage]]:
    """Extract content from a local file.  Dispatches by file extension.

    Uses MarkItDown as the default converter for most formats.
    Falls back to legacy extractors if MarkItDown fails or is disabled.
    Built-in extractors are kept for .md and .txt (already text).

    Returns ``(ArticleInfo, sections, images)`` where *images* is a list
    of ``ExtractedImage`` objects (non-empty only for formats that support
    image extraction: .docx, .pptx with legacy extractors).
    Raises ``LocalFileError`` on failure.
    """
    p = pathlib.Path(file_path)
    if not p.exists():
        raise LocalFileError(f"File not found: {file_path}")

    suffix = p.suffix.lower()

    # .md and .txt — built-in extractors (already text, no conversion needed)
    if suffix == ".md":
        info, sections = _extract_markdown(p, config)
        return info, sections, []
    if suffix == ".txt":
        info, sections = _extract_txt(p, config)
        return info, sections, []

    # MHTML/MHT — always use legacy MIME decoder (MarkItDown can't decode MIME)
    if suffix in (".mhtml", ".mht"):
        info, sections = _extract_mhtml(p, config)
        return info, sections, []

    # Image-bearing Office formats: when image extraction is on, route
    # directly to legacy extractors — they produce ExtractedImage objects
    # with position markers for the vision pipeline. MarkItDown's text
    # output for .pptx/.docx drops the image binaries, so the vision step
    # would silently skip every embedded picture.
    if extract_images and suffix == ".pptx":
        return _extract_pptx(p, config, extract_images=True)
    if extract_images and suffix == ".docx":
        return _extract_docx(p, config, extract_images=True)

    # All other formats — try MarkItDown first
    use_markitdown = markitdown_config is not None and markitdown_config.enabled
    if use_markitdown:
        try:
            from .markitdown_bridge import convert_file, MarkItDownError
            # Need full Config for markitdown_bridge, build a minimal proxy
            # by passing markitdown_config through a namespace
            import types
            proxy_config = types.SimpleNamespace(markitdown=markitdown_config)
            info, sections = convert_file(str(p), proxy_config)
            if sections:
                print(f"  [local_file] MarkItDown: {p.name} → "
                      f"{len(sections)} sections, {info.word_count} words", flush=True)
                return info, sections, []
        except Exception as e:
            print(f"  [local_file] MarkItDown failed for {p.name}: {e}", flush=True)
            print(f"  [local_file] Falling back to legacy extractor...", flush=True)

    # Legacy extractors (fallback or when MarkItDown disabled)
    # Formats that support image extraction (images disabled path falls here too)
    if suffix == ".pptx":
        return _extract_pptx(p, config, extract_images=extract_images)
    if suffix == ".docx":
        return _extract_docx(p, config, extract_images=extract_images)

    # Formats without image extraction
    legacy_extractors = {
        ".doc": _extract_doc,
        ".html": _extract_html,
        ".htm": _extract_html,
        ".mhtml": _extract_mhtml,
        ".mht": _extract_mhtml,
        ".ppt": _extract_ppt,
    }
    extractor = legacy_extractors.get(suffix)
    if extractor:
        info, sections = extractor(p, config)
        return info, sections, []

    # New formats (only via MarkItDown, no legacy extractor)
    if suffix in (".xlsx", ".xls", ".csv", ".json", ".xml", ".epub", ".msg", ".zip"):
        if not use_markitdown:
            raise LocalFileError(
                f"Format {suffix} requires MarkItDown. "
                f"Set markitdown.enabled: true in config.yaml")
        # MarkItDown already tried and failed above
        raise LocalFileError(f"Failed to extract {suffix} file: {p.name}")

    raise LocalFileError(f"Unsupported file format: {suffix}")


# ---------------------------------------------------------------------------
# Shared helpers
# ---------------------------------------------------------------------------

def _guess_title(sections: List[ArticleSection]) -> str:
    """Derive title from the first heading, if any."""
    for sec in sections:
        if sec.heading:
            return sec.heading
    return ""


def _file_date(p: pathlib.Path) -> str:
    """Return file modification date as YYYY-MM-DD."""
    return datetime.fromtimestamp(p.stat().st_mtime).strftime("%Y-%m-%d")


def _build_info(
    path: pathlib.Path,
    sections: List[ArticleSection],
    *,
    title: str = "",
    author: str = "",
    publish_date: str = "",
    language: Optional[str] = None,
    description: str = "",
) -> ArticleInfo:
    """Build an ``ArticleInfo`` from file metadata and parsed sections."""
    body = "\n\n".join(
        (sec.heading + "\n" + sec.body if sec.heading else sec.body)
        for sec in sections
    )
    word_count = len(body.split())
    return ArticleInfo(
        title=title or _guess_title(sections) or path.stem,
        url=path.resolve().as_uri(),
        author=author,
        site_name="",
        publish_date=publish_date or _file_date(path),
        language=language,
        description=description,
        word_count=word_count,
        sections=sections,
    )


# ---------------------------------------------------------------------------
# .md  — Markdown files
# ---------------------------------------------------------------------------

def _extract_markdown(
    path: pathlib.Path,
    config: "LocalFilesConfig",
) -> Tuple[ArticleInfo, List[ArticleSection]]:
    text = path.read_text(encoding="utf-8", errors="replace")

    # Parse YAML frontmatter (between --- markers)
    frontmatter: dict = {}
    body = text
    if text.startswith("---"):
        parts = text.split("---", 2)
        if len(parts) >= 3:
            try:
                import yaml
                frontmatter = yaml.safe_load(parts[1]) or {}
            except Exception:
                pass
            body = parts[2]

    # Reuse the heading parser from pdf.py (skip PDF-specific cleaning)
    from .pdf import parse_markdown_to_sections
    sections = parse_markdown_to_sections(body, pdf_cleanup=False)

    if not sections:
        sections = [ArticleSection(heading="", level=2, body=body.strip())]

    info = _build_info(
        path, sections,
        title=str(frontmatter.get("title", "")),
        author=str(frontmatter.get("author", "")),
        publish_date=str(frontmatter.get("date", "")),
        language=frontmatter.get("language") or frontmatter.get("lang"),
        description=str(frontmatter.get("description", "")),
    )
    return info, sections


# ---------------------------------------------------------------------------
# .txt  — Plain text files
# ---------------------------------------------------------------------------

_RE_ALLCAPS = re.compile(r"^[A-Z][A-Z0-9 :,.\-]{2,}$")
_RE_COLON_HEADING = re.compile(r"^(.{3,60}):\s*$")


def _detect_txt_headings(text: str) -> List[ArticleSection]:
    """Split text into sections by detecting pseudo-headings."""
    lines = text.split("\n")
    sections: List[ArticleSection] = []
    current_heading = ""
    current_lines: List[str] = []

    def _flush():
        body = "\n".join(current_lines).strip()
        if body or current_heading:
            sections.append(ArticleSection(
                heading=current_heading,
                level=2,
                body=body,
            ))

    for line in lines:
        stripped = line.strip()
        is_heading = False
        if stripped and _RE_ALLCAPS.match(stripped):
            is_heading = True
        elif stripped and _RE_COLON_HEADING.match(stripped):
            is_heading = True

        if is_heading:
            _flush()
            current_heading = stripped.rstrip(":")
            current_lines = []
        else:
            current_lines.append(line)

    _flush()
    return sections


def _extract_txt(
    path: pathlib.Path,
    config: "LocalFilesConfig",
) -> Tuple[ArticleInfo, List[ArticleSection]]:
    text = path.read_text(encoding="utf-8", errors="replace")

    if config.detect_txt_headings:
        sections = _detect_txt_headings(text)
    else:
        sections = []

    # Fallback: single section split on double newlines as paragraphs
    if not sections:
        body = re.sub(r"\n{3,}", "\n\n", text.strip())
        sections = [ArticleSection(heading="", level=2, body=body)]

    info = _build_info(path, sections)
    return info, sections


# ---------------------------------------------------------------------------
# .docx  — Word documents (python-docx)
# ---------------------------------------------------------------------------

def _extract_docx(
    path: pathlib.Path,
    config: "LocalFilesConfig",
    extract_images: bool = False,
) -> Tuple[ArticleInfo, List[ArticleSection], List[ExtractedImage]]:
    from .deps import ensure_python_docx
    ensure_python_docx()
    import docx

    try:
        doc = docx.Document(str(path))
    except Exception as exc:
        raise LocalFileError(f"Failed to open .docx file: {exc}") from exc

    sections: List[ArticleSection] = []
    images: List[ExtractedImage] = []
    current_heading = ""
    current_level = 2
    current_paragraphs: List[str] = []

    def _flush():
        nonlocal current_heading, current_paragraphs
        body = "\n\n".join(current_paragraphs)
        if body or current_heading:
            sections.append(ArticleSection(
                heading=current_heading,
                level=current_level,
                body=body,
            ))
        current_heading = ""
        current_paragraphs = []

    for para in doc.paragraphs:
        style_name = (para.style.name or "").lower()
        if style_name.startswith("heading"):
            _flush()
            current_heading = para.text.strip()
            level_match = re.search(r"(\d+)", style_name)
            current_level = int(level_match.group(1)) if level_match else 2
            current_paragraphs = []
        else:
            text = para.text.strip()
            if text:
                current_paragraphs.append(text)

    _flush()

    # Extract inline images from document
    if extract_images:
        from .vision import make_image_marker
        try:
            from docx.enum.shape import WD_INLINE_SHAPE_TYPE
            for shape in doc.inline_shapes:
                try:
                    if shape.type == WD_INLINE_SHAPE_TYPE.INLINE_PICTURE:
                        marker = make_image_marker()
                        ct = shape.image.content_type or "image/png"
                        fmt = ct.split("/")[-1]
                        if fmt not in ("png", "jpeg", "jpg", "gif", "webp"):
                            fmt = "png"
                        images.append(ExtractedImage(
                            image_bytes=shape.image.blob,
                            format=fmt,
                            source_label=path.name,
                            position_marker=marker,
                            alt_text="",
                        ))
                        # Append marker to the last section's body
                        if sections:
                            sections[-1] = ArticleSection(
                                heading=sections[-1].heading,
                                level=sections[-1].level,
                                body=sections[-1].body + "\n\n" + marker,
                            )
                except Exception:
                    pass  # Skip unreadable shapes
        except ImportError:
            pass  # WD_INLINE_SHAPE_TYPE not available in older python-docx

    # Optionally include tables
    if config.include_tables and doc.tables:
        for table in doc.tables:
            rows = []
            for row in table.rows:
                cells = [cell.text.strip() for cell in row.cells]
                rows.append(" | ".join(cells))
            if rows:
                # Markdown table with header separator
                table_md = rows[0] + "\n" + " | ".join(["---"] * len(table.rows[0].cells))
                if len(rows) > 1:
                    table_md += "\n" + "\n".join(rows[1:])
                sections.append(ArticleSection(
                    heading="",
                    level=2,
                    body=table_md,
                ))

    if not sections:
        raise LocalFileError("No content found in .docx file.")

    # Extract metadata from document properties
    props = doc.core_properties
    title = props.title or ""
    author = props.author or ""
    date_val = props.created or props.modified
    publish_date = date_val.strftime("%Y-%m-%d") if date_val else ""

    info = _build_info(
        path, sections,
        title=title,
        author=author,
        publish_date=publish_date,
    )
    return info, sections, images


# ---------------------------------------------------------------------------
# .doc  — Legacy Word (mammoth → HTML → trafilatura)
# ---------------------------------------------------------------------------

def _extract_doc(
    path: pathlib.Path,
    config: "LocalFilesConfig",
) -> Tuple[ArticleInfo, List[ArticleSection]]:
    from .deps import ensure_mammoth
    ensure_mammoth()
    import mammoth

    try:
        with open(path, "rb") as f:
            result = mammoth.convert_to_html(f)
        html = result.value
    except Exception as exc:
        raise LocalFileError(f"Failed to convert .doc file: {exc}") from exc

    if not html or not html.strip():
        raise LocalFileError("No content extracted from .doc file.")

    # Reuse the article extraction pipeline for the converted HTML
    from .article import extract_article
    from .config import ArticlesConfig

    articles_config = ArticlesConfig(
        include_tables=config.include_tables,
        min_content_length=config.min_content_length,
    )
    info, sections, _imgs = extract_article(html, path.resolve().as_uri(), articles_config)

    # Override with file-based metadata
    if not info.title or info.title == "Untitled":
        info.title = _guess_title(sections) or path.stem
    info.publish_date = info.publish_date if info.publish_date != "unknown" else _file_date(path)

    return info, sections


# ---------------------------------------------------------------------------
# .html / .htm  — Local HTML files (trafilatura)
# ---------------------------------------------------------------------------

def _extract_html(
    path: pathlib.Path,
    config: "LocalFilesConfig",
) -> Tuple[ArticleInfo, List[ArticleSection]]:
    html = path.read_text(encoding="utf-8", errors="replace")

    from .article import extract_article
    from .config import ArticlesConfig

    articles_config = ArticlesConfig(
        include_tables=config.include_tables,
        min_content_length=config.min_content_length,
    )
    info, sections, _imgs = extract_article(html, path.resolve().as_uri(), articles_config)

    # Override with file-based metadata
    if not info.title or info.title == "Untitled":
        info.title = _guess_title(sections) or path.stem
    info.publish_date = info.publish_date if info.publish_date != "unknown" else _file_date(path)

    return info, sections


# ---------------------------------------------------------------------------
# .mhtml / .mht  — MHTML web archive (MIME-encoded HTML)
# ---------------------------------------------------------------------------

def _decode_mhtml(path: pathlib.Path) -> Tuple[str, str]:
    """Parse an MHTML archive and return ``(html, subject)``.

    Uses Python's stdlib ``email`` module to walk MIME parts and extract
    the first ``text/html`` part.  Handles quoted-printable and base64
    Content-Transfer-Encoding transparently.

    The *subject* comes from the MIME ``Subject`` header — browsers embed
    the page title there when saving as MHTML.  Returns empty string when
    the header is absent.
    """
    raw = path.read_bytes()
    msg = email.message_from_bytes(raw, policy=email.policy.default)
    subject = str(msg.get("Subject", "") or "")

    for part in msg.walk():
        if part.get_content_type() == "text/html":
            try:
                html = part.get_content()
            except (KeyError, LookupError, UnicodeDecodeError):
                payload = part.get_payload(decode=True)
                if payload is None:
                    continue
                charset = part.get_content_charset() or "utf-8"
                try:
                    html = payload.decode(charset, errors="replace")
                except (LookupError, UnicodeDecodeError):
                    html = payload.decode("utf-8", errors="replace")
            if html and html.strip():
                return html, subject

    raise LocalFileError(f"No text/html part found in MHTML archive: {path.name}")


def _extract_mhtml(
    path: pathlib.Path,
    config: "LocalFilesConfig",
) -> Tuple[ArticleInfo, List[ArticleSection]]:
    html, subject = _decode_mhtml(path)

    from .article import extract_article
    from .config import ArticlesConfig

    articles_config = ArticlesConfig(
        include_tables=config.include_tables,
        min_content_length=config.min_content_length,
    )
    info, sections, _imgs = extract_article(html, path.resolve().as_uri(), articles_config)

    # Title priority: MIME Subject > trafilatura > first heading > filename
    # Subject is most reliable — browsers embed the <title> there when saving.
    if subject:
        info.title = subject
    elif not info.title or info.title == "Untitled":
        info.title = _guess_title(sections) or path.stem
    info.publish_date = info.publish_date if info.publish_date != "unknown" else _file_date(path)

    return info, sections


# ---------------------------------------------------------------------------
# .pptx  — PowerPoint presentations (python-pptx)
# ---------------------------------------------------------------------------

def _extract_pptx(
    path: pathlib.Path,
    config: "LocalFilesConfig",
    extract_images: bool = False,
) -> Tuple[ArticleInfo, List[ArticleSection], List[ExtractedImage]]:
    from .deps import ensure_python_pptx
    ensure_python_pptx()
    from pptx import Presentation

    try:
        prs = Presentation(str(path))
    except Exception as exc:
        raise LocalFileError(f"Failed to open .pptx file: {exc}") from exc

    sections: List[ArticleSection] = []
    images: List[ExtractedImage] = []

    for slide_num, slide in enumerate(prs.slides, 1):
        # Extract slide title
        title = ""
        title_shape = slide.shapes.title
        if title_shape and title_shape.has_text_frame:
            title = title_shape.text_frame.text.strip()

        # Extract body text from all non-title text frames
        body_parts: List[str] = []
        for shape in slide.shapes:
            if shape == title_shape:
                continue
            if shape.has_text_frame:
                for para in shape.text_frame.paragraphs:
                    text = para.text.strip()
                    if text:
                        if para.level > 0:
                            body_parts.append(f"{'  ' * para.level}- {text}")
                        else:
                            body_parts.append(f"- {text}")

        # Extract images from picture shapes
        if extract_images:
            from pptx.enum.shapes import MSO_SHAPE_TYPE
            from .vision import make_image_marker
            for shape in slide.shapes:
                if shape == title_shape:
                    continue
                try:
                    if shape.shape_type == MSO_SHAPE_TYPE.PICTURE:
                        img = shape.image
                        marker = make_image_marker()
                        ct = img.content_type or "image/png"
                        fmt = ct.split("/")[-1]
                        if fmt == "jpeg":
                            fmt = "jpeg"
                        elif fmt not in ("png", "gif", "webp"):
                            fmt = "png"
                        images.append(ExtractedImage(
                            image_bytes=img.blob,
                            format=fmt,
                            source_label=f"Slide {slide_num}",
                            position_marker=marker,
                            alt_text=shape.name or "",
                        ))
                        body_parts.append(marker)
                except Exception:
                    pass  # Skip shapes that can't be read as images

        # Extract tables inline
        if config.include_tables:
            for shape in slide.shapes:
                if shape.has_table:
                    table = shape.table
                    rows = []
                    for row in table.rows:
                        cells = [cell.text.strip() for cell in row.cells]
                        rows.append(" | ".join(cells))
                    if rows:
                        table_md = rows[0] + "\n" + " | ".join(
                            ["---"] * len(table.columns))
                        if len(rows) > 1:
                            table_md += "\n" + "\n".join(rows[1:])
                        body_parts.append(table_md)

        # Extract speaker notes
        notes_text = ""
        if config.include_speaker_notes and slide.has_notes_slide:
            notes_text = slide.notes_slide.notes_text_frame.text.strip()

        # Compose section body
        body = "\n\n".join(body_parts)
        if notes_text:
            notes_quoted = notes_text.replace("\n", "\n> ")
            body += f"\n\n> **Speaker Notes:** {notes_quoted}"

        heading = title or f"Slide {slide_num}"

        # Skip completely empty slides
        if not body and not title:
            continue

        sections.append(ArticleSection(
            heading=heading,
            level=2,
            body=body,
        ))

    if not sections:
        raise LocalFileError("No content found in .pptx file.")

    # Extract metadata from presentation properties
    props = prs.core_properties
    prs_title = props.title or ""
    author = props.author or ""
    date_val = props.created or props.modified
    publish_date = date_val.strftime("%Y-%m-%d") if date_val else ""

    info = _build_info(
        path, sections,
        title=prs_title,
        author=author,
        publish_date=publish_date,
    )
    return info, sections, images


# ---------------------------------------------------------------------------
# .ppt  — Legacy PowerPoint (not supported)
# ---------------------------------------------------------------------------

def _extract_ppt(
    path: pathlib.Path,
    config: "LocalFilesConfig",
) -> Tuple[ArticleInfo, List[ArticleSection]]:
    raise LocalFileError(
        "Legacy .ppt format is not directly supported. "
        "Please convert to .pptx using PowerPoint or LibreOffice, then retry."
    )
